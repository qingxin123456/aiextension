from flask_cors import CORS
from flask import Flask, request, jsonify
from sparkai.llm.llm import ChatSparkLLM, ChunkPrintHandler
from sparkai.core.messages import ChatMessage
import re
import hashlib
from collections import defaultdict
from datetime import datetime
from flask import send_file
from werkzeug.utils import secure_filename
import tempfile
import os
from docx import Document
from pptx import Presentation


code_history = defaultdict(list)
app = Flask(__name__)
CORS(app)
SPARKAI_URL = 'wss://spark-api.xf-yun.com/v1.1/chat'
SPARKAI_APP_ID = 'c5679d7b'
SPARKAI_API_SECRET = 'MGYyOWE4ZTcxODZiODAwYzhkNDQ1Yzdl'
SPARKAI_API_KEY = 'defb43d3e239f8fcbc0aab9a49f493d8'
SPARKAI_DOMAIN = 'lite'
model = ChatSparkLLM(
    spark_api_url="wss://spark-api.xf-yun.com/v1.1/chat",
    spark_app_id="c5679d7b",  # 替换为你的APP ID
    spark_api_secret="MGYyOWE4ZTcxODZiODAwYzhkNDQ1Yzdl",  # 替换为你的API SECRET
    spark_api_key="defb43d3e239f8fcbc0aab9a49f493d8",
    spark_llm_domain="lite",
    streaming=False
)
# 配置上传参数
UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'txt', 'c', 'java', 'py'}
MAX_CONTENT_LENGTH = 10 * 1024 * 1024  # 10MB

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = MAX_CONTENT_LENGTH

# 编程领域专用配置
CODE_KEYWORDS = {
    'debug': ['错误', '异常', 'bug', '调试', 'traceback'],
    'review': ['优化', '改进', '重构', '规范', '可读性'],
    'implement': ['实现', '编写', '算法', '函数', '类', '模块'],
    'explain': ['解释', '什么', '如何工作', '原理']

}

PROMPT_TEMPLATES = {
    'general': {
        'instruction': "作为全栈工程师，请从以下角度进行技术分析：\n1. 核心功能\n2. 实现逻辑\n3. 改进建议\n4. 扩展方案\n技术问题：",
        'examples': [
            "通用分析：该代码模块实现了...",
            "典型建议：建议增加日志记录和异常处理"
        ]
    },
    'debug': {
        'instruction': "作为资深调试专家，请按以下步骤分析：\n1. 定位异常位置\n2. 解释错误原因\n3. 提供修复方案\n4. 建议预防措施\n问题描述：",
        'examples': [
            "类似错误案例：当遇到NullPointerException时，通常是因为...",
            "常见解决方案：添加空值检查，使用Optional包装对象"
        ]
    },
    'review': {
        'instruction': "作为代码评审专家，请从以下维度分析：\n- 代码规范\n- 性能优化\n- 异常处理\n- 可维护性\n- 安全风险\n代码片段：",
        'examples': [
            "优化建议：使用设计模式重构冗余代码",
            "安全提示：注意防范SQL注入漏洞"
        ]
    },
    'implement': {
        'instruction': "作为系统架构师，请：\n1. 分析需求\n2. 设计实现方案\n3. 编写示例代码\n4. 说明注意事项\n功能需求：",
        'examples': [
            "示例实现：使用Dijkstra算法解决最短路径问题",
            "性能考量：时间复杂度O(n^2)可优化为O(nlogn)"
        ]
    },
    'explain': {
        'instruction': "作为技术导师，请解释：\n1. 工作原理\n2. 执行流程\n3. 应用场景\n4. 最佳实践\n技术问题：",
        'examples': [
            "内存管理：Python使用引用计数和垃圾回收机制",
            "异步编程：async/await的工作原理"
        ]
    }
}

CACHE = defaultdict(str)

spark = ChatSparkLLM(
    spark_api_url=SPARKAI_URL,
    spark_app_id=SPARKAI_APP_ID,
    spark_api_key=SPARKAI_API_KEY,
    spark_api_secret=SPARKAI_API_SECRET,
    spark_llm_domain=SPARKAI_DOMAIN,
    streaming=False,
)


def detect_code_intent(text):
    """代码意图检测"""
    text = text.lower()
    for intent, keywords in CODE_KEYWORDS.items():
        if any(kw in text for kw in keywords):
            return intent
    return 'general'

def analyze_code_content(content):
    """基础代码分析函数"""
    return {
        'line_count': len(content.split('\n')),
        'has_functions': 'def ' in content or 'function ' in content,
        'error_patterns': re.findall(r'except|error', content, re.IGNORECASE)
    }



def extract_file_content(file_path, filename):
    """根据文件类型提取内容"""
    ext = filename.rsplit('.', 1)[1].lower()

    # 处理代码文件
    if ext in ALLOWED_EXTENSIONS['code'] + ALLOWED_EXTENSIONS['text']:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read(1024 * 1024)  # 最多读取1MB文本

    # 处理Word文档
    if ext in ['doc', 'docx']:
        doc = Document(file_path)
        return '\n'.join([para.text for para in doc.paragraphs])

    # 处理PPT文档
    if ext in ['ppt', 'pptx']:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)

    return "Unsupported file type"


# 创建上传目录
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)


def allowed_file(filename):
    return '.' in filename and \
        filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


@app.route('/analyze', methods=['POST'])
def analyze_file():
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files['file']
    content = file.read().decode('utf-8')[:5000]  # 限制分析前5000字符

    # 构造对话消息
    messages = [
        ChatMessage(role="user",
                    content=f"请分析以下代码文件：\n```\n{content}\n```\n分析要求：\n1.核心功能\n2.代码结构\n3.改进建议")
    ]

    # 调用大模型
    try:
        response = model.generate([messages])  # 正确调用方法
        analysis = response.generations[0][0].text

        return jsonify({
            "filename": file.filename,
            "analysis": analysis
        })

    except Exception as e:
        return jsonify({"error": str(e)}), 500

def enhance_code_prompt(prompt, intent):
    """代码领域提示工程"""
    template = PROMPT_TEMPLATES.get(intent, PROMPT_TEMPLATES['general'])
    enhanced = template['instruction']

    # 添加示例增强
    if intent != 'general':
        enhanced += "\n".join(template['examples']) + "\n"

    # 代码结构化处理
    if '```' in prompt:
        enhanced += "\n（识别到代码块，将进行深度分析）\n"

    # 添加分析维度标签
    tags = {
        'debug': "[异常类型][根因分析][修复方案]",
        'review': "[代码异味][优化建议][重构方案]",
        'implement': "[架构设计][算法选择][代码示例]",
        'explain': "[核心原理][执行流程][应用场景]"
    }
    enhanced += f"\n{tags.get(intent, '[技术分析][实现方案][最佳实践]')}\n"

    return enhanced + prompt


def code_chunk_processing(content):
    """代码敏感型分块"""
    code_blocks = re.split(r'(```.*?\n|\n\s*\n)', content, flags=re.DOTALL)
    chunks = []
    current_chunk = []
    current_length = 0

    for block in code_blocks:
        block_length = len(block)
        if current_length + block_length > 1500:
            chunks.append("".join(current_chunk))
            current_chunk = [block]
            current_length = block_length
        else:
            current_chunk.append(block)
            current_length += block_length

    if current_chunk:
        chunks.append("".join(current_chunk))
    return chunks


@app.route('/query', methods=['POST'])
def query_sparkai():
    data = request.json
    user_input = data.get('query')

    # 生成缓存键（包含代码特征）
    cache_key = hashlib.md5(user_input.encode()).hexdigest()

    if CACHE[cache_key]:
        return jsonify({"response": CACHE[cache_key]})

    intent = detect_code_intent(user_input)
    enhanced_prompt = enhance_code_prompt(user_input, intent)

    # 多角色分析
    roles = {
        'debug': "资深调试专家",
        'review': "首席代码审查员",
        'implement': "系统架构师",
        'explain': "技术导师"
    }
    role = roles.get(intent, "高级工程师")

    messages = [
        ChatMessage(role="system", content=f"你现在是{role}，请用专业但易懂的方式回答"),
        ChatMessage(role="user", content=enhanced_prompt)
    ]

    handler = ChunkPrintHandler()
    response = spark.generate([messages], callbacks=[handler])

    if response and response.generations:
        result = response.generations[0][0].text
        # 后处理格式化
        if '代码示例' in result:
            result = re.sub(r'(```)(?!python)', r'\1python', result)
        CACHE[cache_key] = result
        return jsonify({"response": result})
    else:
        return jsonify({"error": "分析失败"}), 500



@app.route('/code/history', methods=['POST'])
def handle_code_history():
        code = request.json.get('code')
        user = request.json.get('user', 'default')

        timestamp = datetime.now().isoformat()
        code_history[user].append({
            "time": timestamp,
            "content": code,
            "hash": hashlib.md5(code.encode()).hexdigest()
        })

        return jsonify({
            "versions": code_history[user][-5:],  # 返回最近5个版本
            "current_hash": code_history[user][-1]['hash']
        })

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)